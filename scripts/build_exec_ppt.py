"""Build oapw executive stakeholder presentation."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Brand palette ─────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)   # deep navy — primary background
COBALT      = RGBColor(0x1A, 0x56, 0xDB)   # bright blue — accent / headings
TEAL        = RGBColor(0x05, 0xB7, 0xA8)   # teal — highlights / callouts
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xFA)
MID_GRAY    = RGBColor(0x8A, 0x9B, 0xB0)
DARK_GRAY   = RGBColor(0x2D, 0x3A, 0x4A)
GREEN       = RGBColor(0x10, 0xB9, 0x81)   # ✓ positive / metric
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)   # caution

# Slide dimensions (widescreen 16:9)
W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def solid_fill(shape, color: RGBColor):
    """Fill a shape with a solid colour."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    solid_fill(shape, color)
    shape.line.fill.background()   # no border
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    """Add a text box with consistent styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_bullet_box(slide, items: list[str], left, top, width, height,
                   font_size=16, color=WHITE, bullet_color=None,
                   line_spacing=1.15, font_name="Calibri"):
    """Add a multi-bullet text box."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)

        # bullet character
        run_bullet = p.add_run()
        run_bullet.text = "▸  "
        run_bullet.font.size = Pt(font_size - 2)
        run_bullet.font.color.rgb = bullet_color or TEAL
        run_bullet.font.name = font_name

        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = font_name

    return txBox


def set_bg(slide, color: RGBColor):
    """Set slide background to a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def accent_bar(slide, color=COBALT, height=Inches(0.06)):
    """Add a thin colour bar at the very top of the slide."""
    add_rect(slide, 0, 0, W, height, color)


def slide_number_label(slide, num: int, total: int):
    """Add a subtle slide number in bottom-right."""
    add_textbox(slide, f"{num} / {total}",
                W - Inches(1.2), H - Inches(0.4),
                Inches(1.0), Inches(0.3),
                font_size=10, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — Cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, NAVY)

    # Left colour stripe
    add_rect(slide, 0, 0, Inches(0.45), H, COBALT)

    # Teal accent bottom bar
    add_rect(slide, 0, H - Inches(0.55), W, Inches(0.55), TEAL)

    # Company / product tag
    add_textbox(slide, "CONFIDENTIAL  ·  EXECUTIVE BRIEFING",
                Inches(0.75), Inches(0.55), Inches(11), Inches(0.5),
                font_size=11, color=TEAL, bold=True)

    # Main title
    add_textbox(slide, "oapw",
                Inches(0.75), Inches(1.4), Inches(11), Inches(1.4),
                font_size=72, bold=True, color=WHITE, font_name="Calibri")

    # Subtitle
    add_textbox(slide, "AI-Powered Quality Engineering Platform",
                Inches(0.75), Inches(2.85), Inches(10), Inches(0.75),
                font_size=28, bold=False, color=LIGHT_GRAY)

    # Tagline
    add_textbox(slide,
                "Local-first · Zero cloud costs · Autonomous regression · "
                "Knows your codebase",
                Inches(0.75), Inches(3.65), Inches(10), Inches(0.6),
                font_size=16, color=TEAL, italic=True)

    # Meta
    add_textbox(slide, "May 2026  ·  Engineering Leadership Review",
                Inches(0.75), H - Inches(0.95), Inches(8), Inches(0.4),
                font_size=13, color=MID_GRAY)


def slide_agenda(prs):
    """Slide 2 — Agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 2, 14)

    add_textbox(slide, "Today's Agenda",
                Inches(0.6), Inches(0.35), Inches(10), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    items = [
        ("01", "The Problem with Traditional QA"),
        ("02", "What oapw Does — and Why It's Different"),
        ("03", "Platform Capabilities at a Glance"),
        ("04", "The Knowledge Base Advantage"),
        ("05", "Autonomous QA Agent"),
        ("06", "Multi-Faceted Verification Suite"),
        ("07", "Developer & CI/CD Experience"),
        ("08", "ROI & Cost Model"),
        ("09", "Delivery Track Record"),
        ("10", "Roadmap & Next Steps"),
    ]

    col1 = items[:5]
    col2 = items[5:]

    for i, (num, label) in enumerate(col1):
        y = Inches(1.5) + i * Inches(0.95)
        add_rect(slide, Inches(0.6), y, Inches(0.55), Inches(0.55), COBALT)
        add_textbox(slide, num, Inches(0.6), y, Inches(0.55), Inches(0.55),
                    font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, Inches(1.3), y + Inches(0.04),
                    Inches(5.2), Inches(0.55),
                    font_size=16, color=LIGHT_GRAY)

    for i, (num, label) in enumerate(col2):
        y = Inches(1.5) + i * Inches(0.95)
        add_rect(slide, Inches(7.1), y, Inches(0.55), Inches(0.55), TEAL)
        add_textbox(slide, num, Inches(7.1), y, Inches(0.55), Inches(0.55),
                    font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, Inches(7.8), y + Inches(0.04),
                    Inches(5.2), Inches(0.55),
                    font_size=16, color=LIGHT_GRAY)


def slide_problem(prs):
    """Slide 3 — The Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 3, 14)

    add_textbox(slide, "The Problem with Traditional QA",
                Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    pains = [
        ("💸", "Cloud AI costs spiral with every CI run",
         "GPT-4 / Claude API charges accumulate fast — $0.01–0.06 per LLM call, "
         "thousands of calls per day in active teams."),
        ("🔒", "Sensitive data leaves the building",
         "Source code, Jira tickets, acceptance criteria, and test data sent to "
         "cloud LLMs violate data residency and IP policies."),
        ("🤖", "AI generates generic, brittle tests",
         "Without context about your app, AI writes 'fill the email field' instead "
         "of tests grounded in your actual acceptance criteria."),
        ("🩹", "Locators break overnight — manually fixed",
         "Every UI refactor silently breaks 10–40% of selectors. Engineers spend "
         "hours updating XPaths instead of building features."),
    ]

    for i, (icon, heading, detail) in enumerate(pains):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(1.4) + row * Inches(2.5)

        add_rect(slide, x, y, Inches(5.9), Inches(2.2), DARK_GRAY)
        add_textbox(slide, icon, x + Inches(0.2), y + Inches(0.15),
                    Inches(0.6), Inches(0.6), font_size=26)
        add_textbox(slide, heading,
                    x + Inches(0.85), y + Inches(0.18),
                    Inches(4.8), Inches(0.55),
                    font_size=15, bold=True, color=AMBER)
        add_textbox(slide, detail,
                    x + Inches(0.2), y + Inches(0.75),
                    Inches(5.5), Inches(1.3),
                    font_size=13, color=LIGHT_GRAY)


def slide_solution(prs):
    """Slide 4 — What oapw Does."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 4, 14)

    add_textbox(slide, "What oapw Does — and Why It's Different",
                Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    # Big statement
    add_rect(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(1.1), COBALT)
    add_textbox(slide,
                '"Tell it what to test. It knows your app. It finds the bugs."',
                Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                italic=True)

    differentiators = [
        ("100% Local", "All AI inference runs on Ollama — your machine, your network. "
         "Zero data egress. Zero per-token billing."),
        ("Context-Aware", "Ingests Jira, Confluence, and your C#/TypeScript source "
         "so generated tests match real acceptance criteria."),
        ("Self-Healing", "Locators auto-repair using fingerprinting + LLM fallback. "
         "UI refactors no longer break the test suite."),
        ("Autonomous Agent", "oapw qa 'regress the login flow' — the agent selects tests, "
         "runs them, judges failures, and drafts Jira bugs."),
    ]

    for i, (heading, detail) in enumerate(differentiators):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(2.65) + row * Inches(2.1)

        add_rect(slide, x, y, Inches(0.18), Inches(1.7), TEAL)
        add_textbox(slide, heading,
                    x + Inches(0.35), y + Inches(0.1),
                    Inches(5.8), Inches(0.5),
                    font_size=16, bold=True, color=TEAL)
        add_textbox(slide, detail,
                    x + Inches(0.35), y + Inches(0.6),
                    Inches(5.8), Inches(1.0),
                    font_size=14, color=LIGHT_GRAY)


def slide_capabilities(prs):
    """Slide 5 — Platform Capabilities."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 5, 14)

    add_textbox(slide, "Platform Capabilities at a Glance",
                Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    phases = [
        ("Core Infrastructure",   "Multi-layer cache · BrowserManager · Config"),
        ("AI Browser Actions",    "AiPage · NL planner · Step executor"),
        ("Self-Healing Locators", "Fingerprints · 3 fallback strategies · SQLite log"),
        ("Knowledge Base",        "ChromaDB · RAG retrieval · Jira boost"),
        ("Hybrid API + UI",       "Shared cookie jar · ApiClient · PII masker"),
        ("Test Generator",        "from-jira · from-story · smoke crawler · edge mutations"),
        ("Agent System",          "AgentRunner · LoopGuard · Hook system · Replan"),
        ("QA Agent Mode",         "GoalParser · TestSelector · JudgmentEngine · Investigator"),
        ("Verification Suite",    "Accessibility (axe) · Performance (Web Vitals) · Visual diff"),
        ("Productionization",     "pytest plugin · oapw init · GitHub Actions · Examples"),
    ]

    for i, (name, detail) in enumerate(phases):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * Inches(6.45)
        y = Inches(1.35) + row * Inches(1.12)

        num_color = COBALT if col == 0 else TEAL
        add_rect(slide, x, y + Inches(0.05), Inches(0.42), Inches(0.42), num_color)
        add_textbox(slide, f"P{i+1}",
                    x, y + Inches(0.05), Inches(0.42), Inches(0.42),
                    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, name,
                    x + Inches(0.52), y + Inches(0.05),
                    Inches(5.6), Inches(0.38),
                    font_size=14, bold=True, color=WHITE)
        add_textbox(slide, detail,
                    x + Inches(0.52), y + Inches(0.45),
                    Inches(5.6), Inches(0.45),
                    font_size=12, color=MID_GRAY)


def slide_kb(prs):
    """Slide 6 — Knowledge Base Advantage."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 6, 14)

    add_textbox(slide, "The Knowledge Base Advantage",
                Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    # Left panel — sources
    add_rect(slide, Inches(0.4), Inches(1.3), Inches(4.0), Inches(5.6), DARK_GRAY)
    add_textbox(slide, "What Gets Indexed",
                Inches(0.6), Inches(1.45), Inches(3.6), Inches(0.5),
                font_size=16, bold=True, color=TEAL)

    sources = [
        "📋  Jira tickets — summaries, AC, status",
        "📄  Confluence — design docs, test plans",
        "⌨️   C# source — classes, methods, XML docs",
        "⚛️   TypeScript/React — components, hooks",
        "🔀  Git history — incremental SHA sync",
    ]
    add_bullet_box(slide, sources, Inches(0.5), Inches(2.0),
                   Inches(3.7), Inches(3.5), font_size=13,
                   color=LIGHT_GRAY, bullet_color=COBALT)

    # Arrow
    add_textbox(slide, "→",
                Inches(4.55), Inches(3.7), Inches(0.6), Inches(0.6),
                font_size=36, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Middle panel — KB
    add_rect(slide, Inches(5.2), Inches(1.3), Inches(3.0), Inches(5.6), COBALT)
    add_textbox(slide, "ChromaDB\nVector Store",
                Inches(5.3), Inches(2.1), Inches(2.8), Inches(1.0),
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "cosine similarity\nnomic-embed-text\nL2 embedding cache",
                Inches(5.3), Inches(3.2), Inches(2.8), Inches(1.2),
                font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "2,500+\ndocs indexed",
                Inches(5.3), Inches(4.55), Inches(2.8), Inches(1.0),
                font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Arrow
    add_textbox(slide, "→",
                Inches(8.3), Inches(3.7), Inches(0.6), Inches(0.6),
                font_size=36, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Right panel — output
    add_rect(slide, Inches(9.0), Inches(1.3), Inches(4.0), Inches(5.6), DARK_GRAY)
    add_textbox(slide, "AI Gets Context",
                Inches(9.2), Inches(1.45), Inches(3.6), Inches(0.5),
                font_size=16, bold=True, color=TEAL)

    outputs = [
        "Tests match real AC, not guesses",
        "Field names from source code",
        "Error messages from design docs",
        "Azure AD flow from Confluence",
        "Jira-linked docs boosted ×1.2",
    ]
    add_bullet_box(slide, outputs, Inches(9.1), Inches(2.0),
                   Inches(3.7), Inches(3.5), font_size=13,
                   color=LIGHT_GRAY, bullet_color=GREEN)


def slide_qa_agent(prs):
    """Slide 7 — Autonomous QA Agent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 7, 14)

    add_textbox(slide, "Autonomous QA Agent",
                Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    # Command box
    add_rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.75), DARK_GRAY)
    add_textbox(slide, '$ oapw qa  "regression of the login flow on QA"',
                Inches(0.8), Inches(1.38), Inches(12), Inches(0.6),
                font_size=17, bold=True, color=GREEN, font_name="Courier New")

    # Pipeline steps
    pipeline = [
        (COBALT, "1\nGoal\nParser",
         "NL → structured intent\nScope · Features\nEnvironment · Jira refs"),
        (COBALT, "2\nTest\nSelector",
         "Memory + KB sources\nScope tier filter\nRelevance ranking"),
        (COBALT, "3\nSmart\nExecutor",
         "pytest JSON report\nor AgentRunner\nparallel execution"),
        (TEAL,   "4\nJudgment\nEngine",
         "real_bug · flaky\nenv_issue · unclear\nconfidence score"),
        (TEAL,   "5\nInvestigator",
         "Jira history lookup\nGit log correlation\nBug report draft"),
        (GREEN,  "6\nConsole\nReporter",
         "Rich table summary\nPass rate · Real bugs\nJira draft ready"),
    ]

    step_w = Inches(1.88)
    for i, (color, label, detail) in enumerate(pipeline):
        x = Inches(0.45) + i * (step_w + Inches(0.12))
        y = Inches(2.35)

        add_rect(slide, x, y, step_w, Inches(1.1), color)
        add_textbox(slide, label, x, y, step_w, Inches(1.1),
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_textbox(slide, detail,
                    x, y + Inches(1.15), step_w, Inches(1.05),
                    font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

        if i < len(pipeline) - 1:
            add_textbox(slide, "›",
                        x + step_w + Inches(0.01), y + Inches(0.35),
                        Inches(0.12), Inches(0.5),
                        font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Outcome box
    add_rect(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(1.0), DARK_GRAY)
    outcomes = [
        "✓  Pass rate reported",
        "✓  Real bugs vs flaky vs env issues — classified automatically",
        "✓  Jira bug draft generated for real failures — one click to file",
    ]
    add_textbox(slide, "   ".join(outcomes),
                Inches(0.7), Inches(5.0), Inches(12), Inches(0.7),
                font_size=13, color=GREEN, bold=False)

    # Hook system note
    add_textbox(slide,
                "Human-in-loop mode available: --interactive pauses on failures "
                "for CONTINUE / RETRY / ABORT / OVERRIDE decisions",
                Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.45),
                font_size=12, color=MID_GRAY, italic=True)


def slide_verification(prs):
    """Slide 8 — Verification Suite."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 8, 14)

    add_textbox(slide, "Multi-Faceted Verification Suite",
                Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    add_textbox(slide,
                "Three verification dimensions — each a one-liner in pytest tests:",
                Inches(0.6), Inches(1.15), Inches(12), Inches(0.45),
                font_size=16, color=LIGHT_GRAY)

    pillars = [
        (COBALT, "♿  Accessibility",
         "WCAG 2.0 AA",
         [
             "Injects axe-core via Playwright",
             "Groups violations: critical / serious",
             "assert_no_critical()  assert_no_serious()",
             "Zero external service needed",
         ],
         "report = await oapw_accessibility.check(page)\nreport.assert_no_critical()"),
        (TEAL, "⚡  Performance",
         "Web Vitals",
         [
             "TTFB · FCP · LCP from browser API",
             "Resource count + bytes breakdown",
             "assert_ttfb_under(ms)  assert_fcp_under(ms)",
             "No external beacon, no cloud account",
         ],
         "metrics = await oapw_performance.capture(page)\nmetrics.assert_fcp_under(2000)"),
        (GREEN, "👁️  Visual Regression",
         "Pixel Diff",
         [
             "Auto-captures baseline on first run",
             "Pillow pixel-diff with change highlight",
             "Optional LLM description of changes",
             "assert_within_threshold()  (default 2%)",
         ],
         "diff = await oapw_visual.compare(page, 'login')\ndiff.assert_within_threshold()"),
    ]

    for i, (color, title, subtitle, bullets, code) in enumerate(pillars):
        x = Inches(0.4) + i * Inches(4.28)
        y = Inches(1.75)

        add_rect(slide, x, y, Inches(3.95), Inches(0.75), color)
        add_textbox(slide, title, x, y, Inches(3.95), Inches(0.45),
                    font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, subtitle, x, y + Inches(0.44), Inches(3.95), Inches(0.3),
                    font_size=12, color=NAVY, align=PP_ALIGN.CENTER, bold=True)

        add_bullet_box(slide, bullets, x, y + Inches(0.85),
                       Inches(3.95), Inches(1.9),
                       font_size=12, color=LIGHT_GRAY, bullet_color=color)

        add_rect(slide, x, y + Inches(2.85), Inches(3.95), Inches(1.05), DARK_GRAY)
        add_textbox(slide, code, x + Inches(0.1), y + Inches(2.9),
                    Inches(3.75), Inches(0.95),
                    font_size=11, color=GREEN, font_name="Courier New")


def slide_dx(prs):
    """Slide 9 — Developer Experience."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 9, 14)

    add_textbox(slide, "Developer & CI/CD Experience",
                Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    # Left — fixtures table
    add_textbox(slide, "10 pytest fixtures — zero config",
                Inches(0.5), Inches(1.25), Inches(6.0), Inches(0.45),
                font_size=16, bold=True, color=TEAL)

    fixtures = [
        ("oapw_page",          "AiPage with .ai() / .ai_assert()"),
        ("oapw_hybrid",        "Browser + API, shared cookie jar"),
        ("oapw_factory",       "Realistic test data — 5 built-in types"),
        ("oapw_accessibility", "WCAG 2.0 AA axe-core audit"),
        ("oapw_performance",   "TTFB · FCP · LCP measurement"),
        ("oapw_visual",        "Pixel-diff screenshot regression"),
        ("oapw_qa_agent",      "Full autonomous QA pipeline"),
    ]

    for i, (name, desc) in enumerate(fixtures):
        y = Inches(1.85) + i * Inches(0.65)
        add_rect(slide, Inches(0.5), y, Inches(2.5), Inches(0.52), DARK_GRAY)
        add_textbox(slide, name, Inches(0.55), y + Inches(0.05),
                    Inches(2.4), Inches(0.42),
                    font_size=12, bold=True, color=GREEN, font_name="Courier New")
        add_textbox(slide, desc, Inches(3.15), y + Inches(0.08),
                    Inches(3.2), Inches(0.42),
                    font_size=12, color=LIGHT_GRAY)

    # Right — CLI + CI
    add_textbox(slide, "CLI & CI",
                Inches(7.0), Inches(1.25), Inches(5.9), Inches(0.45),
                font_size=16, bold=True, color=TEAL)

    cli_items = [
        "oapw init              scaffold new project in seconds",
        "oapw doctor            verify every dependency",
        "oapw generate from-jira AUTH-42    test from ticket",
        "oapw qa 'smoke test checkout'      autonomous run",
        "oapw run goal '...' --interactive  human-in-loop",
    ]
    for i, cmd in enumerate(cli_items):
        y = Inches(1.85) + i * Inches(0.62)
        add_rect(slide, Inches(7.0), y, Inches(5.9), Inches(0.52), DARK_GRAY)
        add_textbox(slide, cmd, Inches(7.1), y + Inches(0.07),
                    Inches(5.7), Inches(0.42),
                    font_size=11, color=GREEN, font_name="Courier New")

    # GitHub Actions
    add_rect(slide, Inches(7.0), Inches(5.0), Inches(5.9), Inches(1.85), DARK_GRAY)
    add_textbox(slide, "GitHub Actions — bundled template",
                Inches(7.1), Inches(5.05), Inches(5.7), Inches(0.4),
                font_size=13, bold=True, color=TEAL)
    ga_items = [
        "unit-tests job — every push, pull model, run 533 tests",
        "integration-tests job — gated on OAPW_RUN_INTEGRATION=true",
        "Allure report upload, 14-day retention",
    ]
    add_bullet_box(slide, ga_items, Inches(7.1), Inches(5.5),
                   Inches(5.7), Inches(1.25),
                   font_size=12, color=LIGHT_GRAY, bullet_color=COBALT)


def slide_roi(prs):
    """Slide 10 — ROI & Cost Model."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 10, 14)

    add_textbox(slide, "ROI & Cost Model",
                Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    # Big numbers
    metrics = [
        (GREEN,  "$0",           "per LLM call in CI\n(vs $0.01–0.06 cloud)"),
        (GREEN,  "100%",         "data stays local\nZero egress / GDPR risk"),
        (COBALT, "~2 min",       "time saved per\nhealed locator (manual fix)"),
        (TEAL,   "10× faster",   "KB-aware test gen\nvs writing from scratch"),
    ]

    for i, (color, big, small) in enumerate(metrics):
        x = Inches(0.4) + i * Inches(3.2)
        add_rect(slide, x, Inches(1.3), Inches(2.95), Inches(2.3), DARK_GRAY)
        add_textbox(slide, big, x, Inches(1.45), Inches(2.95), Inches(1.15),
                    font_size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, small, x, Inches(2.65), Inches(2.95), Inches(0.8),
                    font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Comparison table
    add_textbox(slide, "Cloud LLM API  vs  oapw (local)",
                Inches(0.6), Inches(3.85), Inches(12), Inches(0.45),
                font_size=16, bold=True, color=WHITE)

    rows = [
        ("", "Cloud LLM (GPT-4 / Claude)", "oapw (Ollama local)"),
        ("Inference cost",    "$0.01–0.06 / call · $500–5k/mo at scale", "$0 after hardware"),
        ("Data privacy",      "Source code & tickets sent externally",    "100% on-premise"),
        ("Latency",           "200–2000 ms per call (network + queue)",   "50–500 ms local"),
        ("Offline capability","Requires internet",                         "Fully offline"),
        ("Vendor lock-in",    "Yes — API keys, model changes",            "None — swap models freely"),
    ]

    col_widths = [Inches(2.1), Inches(5.0), Inches(4.8)]
    col_x = [Inches(0.5), Inches(2.7), Inches(7.8)]

    for r, row in enumerate(rows):
        bg = DARK_GRAY if r % 2 == 0 else RGBColor(0x22, 0x30, 0x45)
        header_row = r == 0
        y = Inches(4.4) + r * Inches(0.46)
        for c, (cell, w, x) in enumerate(zip(row, col_widths, col_x)):
            add_rect(slide, x, y, w, Inches(0.44),
                     COBALT if header_row else bg)
            fc = WHITE if header_row else (LIGHT_GRAY if c > 0 else TEAL)
            add_textbox(slide, cell, x + Inches(0.08), y + Inches(0.05),
                        w - Inches(0.1), Inches(0.38),
                        font_size=12, bold=header_row, color=fc)


def slide_delivery(prs):
    """Slide 11 — Delivery Track Record."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 11, 14)

    add_textbox(slide, "Delivery Track Record",
                Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    # Stats row
    stats = [
        (GREEN,  "10",    "Phases\nDelivered"),
        (COBALT, "533",   "Tests\nPassing"),
        (TEAL,   "50+",   "Source\nModules"),
        (GREEN,  "0",     "Regressions\nIntroduced"),
        (COBALT, "< 12s", "Full Suite\nRun Time"),
    ]

    for i, (color, big, label) in enumerate(stats):
        x = Inches(0.4) + i * Inches(2.55)
        add_rect(slide, x, Inches(1.3), Inches(2.3), Inches(1.6), DARK_GRAY)
        add_textbox(slide, big, x, Inches(1.4), Inches(2.3), Inches(0.85),
                    font_size=40, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, x, Inches(2.25), Inches(2.3), Inches(0.6),
                    font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Phase timeline
    add_textbox(slide, "Phase-by-Phase Delivery",
                Inches(0.6), Inches(3.15), Inches(12), Inches(0.4),
                font_size=15, bold=True, color=TEAL)

    phases = [
        ("P1", "Core Infra"),
        ("P2", "AI Actions"),
        ("P3", "Self-Heal"),
        ("P4", "KB + Jira"),
        ("P5", "Hybrid"),
        ("P6", "Generator"),
        ("P7", "Agent"),
        ("P8", "QA Agent"),
        ("P9", "Verify"),
        ("P10", "Prod"),
    ]

    for i, (num, name) in enumerate(phases):
        x = Inches(0.4) + i * Inches(1.27)
        add_rect(slide, x, Inches(3.7), Inches(1.17), Inches(0.55), COBALT)
        add_textbox(slide, num, x, Inches(3.72), Inches(1.17), Inches(0.4),
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, name, x, Inches(4.28), Inches(1.17), Inches(0.4),
                    font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)
        if i < len(phases) - 1:
            add_textbox(slide, "›", x + Inches(1.17), Inches(3.85),
                        Inches(0.1), Inches(0.3),
                        font_size=14, color=TEAL, align=PP_ALIGN.CENTER)

    # Key quality callouts
    callouts = [
        ("All phases tested before commit",
         "Every phase includes a unit test suite. Phase N cannot ship if Phase N-1 tests regress."),
        ("Atomic commits with docs",
         "Each phase: implementation + CHANGELOG + README + architecture + CLI reference updated atomically."),
        ("Deterministic test isolation",
         "Tests use monkeypatched env vars, in-memory SQLite, and ephemeral ChromaDB — no shared state."),
    ]

    for i, (heading, detail) in enumerate(callouts):
        y = Inches(4.85) + i * Inches(0.83)
        add_rect(slide, Inches(0.18), y, Inches(0.18), Inches(0.65), TEAL)
        add_textbox(slide, heading, Inches(0.5), y + Inches(0.04),
                    Inches(12), Inches(0.38),
                    font_size=14, bold=True, color=WHITE)
        add_textbox(slide, detail, Inches(0.5), y + Inches(0.42),
                    Inches(12), Inches(0.35),
                    font_size=12, color=MID_GRAY)


def slide_roadmap(prs):
    """Slide 12 — Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 12, 14)

    add_textbox(slide, "Roadmap & Next Steps",
                Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    horizons = [
        (GREEN,  "Now — Ready to Use",
         [
             "All 10 phases shipped and tested",
             "pytest plugin: drop-in fixtures for any project",
             "oapw init: new project in < 2 minutes",
             "GitHub Actions CI template included",
             "Full docs: setup, CLI, architecture, KB guide",
         ]),
        (COBALT, "Q3 2026 — Team Rollout",
         [
             "Pilot: 2–3 squads adopt oapw_page + oapw_factory",
             "KB sync job scheduled nightly (Jira + Confluence)",
             "Slack bot: oapw qa results posted to #qa-alerts",
             "Allure dashboard for historical pass-rate trends",
             "Training sessions for QA engineers",
         ]),
        (TEAL,   "Q4 2026 — Scale & Intelligence",
         [
             "MkDocs site for internal docs portal",
             "Parallel code repo sync (unblocked on RAM constraint)",
             "Multi-repo KB sync via repos.yml config file",
             "Proactive bug detection: nightly regression without triggers",
             "Integration with internal Jira to auto-file detected bugs",
         ]),
    ]

    for i, (color, label, items) in enumerate(horizons):
        x = Inches(0.4) + i * Inches(4.28)
        add_rect(slide, x, Inches(1.3), Inches(3.95), Inches(0.6), color)
        add_textbox(slide, label, x, Inches(1.34), Inches(3.95), Inches(0.52),
                    font_size=14, bold=True,
                    color=NAVY if color == GREEN else WHITE,
                    align=PP_ALIGN.CENTER)
        add_bullet_box(slide, items, x, Inches(2.0),
                       Inches(3.95), Inches(4.6),
                       font_size=13, color=LIGHT_GRAY, bullet_color=color)


def slide_decision(prs):
    """Slide 13 — Decision / Ask."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent_bar(slide)
    slide_number_label(slide, 13, 14)

    add_textbox(slide, "What We're Asking For",
                Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    asks = [
        (COBALT, "1", "Approval to pilot with 2–3 squads",
         "Select initial teams. 2-week integration sprint. "
         "We provide onboarding, setup guide, and hands-on support."),
        (COBALT, "2", "Hardware: 1 × 16 GB shared CI runner",
         "Ollama needs 16 GB RAM for qwen2.5:7b quality. "
         "Can share with existing CI infra — no dedicated server required."),
        (TEAL,   "3", "Atlassian API token (read-only)",
         "One team-level token for Jira + Confluence read access. "
         "All data stays on-premise — token used only for KB sync."),
        (TEAL,   "4", "Nightly CI slot for KB sync job",
         "5-minute nightly sync keeps the Knowledge Base current "
         "with sprint changes, new tickets, and updated Confluence pages."),
    ]

    for i, (color, num, heading, detail) in enumerate(asks):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(1.45) + row * Inches(2.3)

        add_rect(slide, x, y, Inches(5.9), Inches(2.1), DARK_GRAY)
        add_rect(slide, x, y, Inches(0.55), Inches(0.55), color)
        add_textbox(slide, num, x, y, Inches(0.55), Inches(0.55),
                    font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, heading, x + Inches(0.7), y + Inches(0.1),
                    Inches(5.0), Inches(0.5),
                    font_size=15, bold=True, color=color)
        add_textbox(slide, detail, x + Inches(0.2), y + Inches(0.7),
                    Inches(5.5), Inches(1.2),
                    font_size=13, color=LIGHT_GRAY)


def slide_close(prs):
    """Slide 14 — Closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    # Full-width teal stripe top
    add_rect(slide, 0, 0, W, Inches(0.55), TEAL)
    # Full-width cobalt stripe bottom
    add_rect(slide, 0, H - Inches(0.55), W, Inches(0.55), COBALT)
    # Left stripe
    add_rect(slide, 0, 0, Inches(0.45), H, COBALT)

    add_textbox(slide, "oapw",
                Inches(0.75), Inches(1.6), Inches(11), Inches(1.2),
                font_size=72, bold=True, color=WHITE, font_name="Calibri")

    add_textbox(slide, "Local AI. Zero Cost. Knows Your App.",
                Inches(0.75), Inches(2.85), Inches(11), Inches(0.7),
                font_size=26, color=TEAL, bold=True)

    add_textbox(slide,
                "533 tests passing  ·  10 phases delivered  ·  "
                "Ready for pilot today",
                Inches(0.75), Inches(3.65), Inches(11), Inches(0.55),
                font_size=17, color=LIGHT_GRAY)

    add_textbox(slide, "Questions & Discussion",
                Inches(0.75), Inches(4.7), Inches(6), Inches(0.55),
                font_size=22, bold=True, color=WHITE)

    # Contact / links area
    add_rect(slide, Inches(0.75), Inches(5.5), Inches(11.5), Inches(0.9), DARK_GRAY)
    add_textbox(slide,
                "Repo: github.com/your-org/ollama-playwright-automation   "
                "·   Docs: docs/setup.md   ·   Demo: oapw doctor",
                Inches(0.95), Inches(5.65), Inches(11.2), Inches(0.55),
                font_size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────

def build(output_path: str):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_agenda(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_capabilities(prs)
    slide_kb(prs)
    slide_qa_agent(prs)
    slide_verification(prs)
    slide_dx(prs)
    slide_roi(prs)
    slide_delivery(prs)
    slide_roadmap(prs)
    slide_decision(prs)
    slide_close(prs)

    prs.save(output_path)
    print(f"✓ Saved: {output_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    import sys, pathlib
    out = sys.argv[1] if len(sys.argv) > 1 else str(
        pathlib.Path(__file__).parent.parent /
        "oapw_executive_presentation.pptx"
    )
    build(out)
